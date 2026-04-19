#!/usr/bin/env python3
"""Generate PPT from Markdown."""

import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def parse_markdown(content: str) -> list[dict]:
    """Parse Markdown into slides."""
    slides = []
    lines = content.strip().split('\n')
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        if line.startswith('# '):
            # Title slide
            if current_slide:
                slides.append(current_slide)
            current_slide = {'type': 'title', 'title': line[2:], 'content': []}
        elif line.startswith('## '):
            # Content slide
            if current_slide:
                slides.append(current_slide)
            current_slide = {'type': 'content', 'title': line[3:], 'content': []}
        elif line.startswith('- '):
            # Bullet point
            if current_slide:
                current_slide['content'].append(line[2:])
    
    if current_slide:
        slides.append(current_slide)
    
    return slides


def create_ppt(slides: list[dict], output_path: str, theme: str = 'default'):
    """Create PPT file from slides."""
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        if slide_data['type'] == 'title':
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            # Add title text box
            left = Inches(0.5)
            top = Inches(3)
            width = Inches(12.33)
            height = Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(44)
            p.font.bold = True
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            # Add title
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12.33)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(32)
            p.font.bold = True
            
            # Add content
            if slide_data['content']:
                left = Inches(0.5)
                top = Inches(1.8)
                width = Inches(12.33)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                
                for i, item in enumerate(slide_data['content']):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(24)
    
    prs.save(output_path)
    print(f"Generated: {output_path} ({len(slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description='Generate PPT from Markdown')
    parser.add_argument('--input', required=True, help='Input Markdown file')
    parser.add_argument('--output', required=True, help='Output PPT file')
    parser.add_argument('--theme', default='default', help='Theme name')
    args = parser.parse_args()
    
    content = Path(args.input).read_text(encoding='utf-8')
    slides = parse_markdown(content)
    create_ppt(slides, args.output, args.theme)


if __name__ == '__main__':
    main()